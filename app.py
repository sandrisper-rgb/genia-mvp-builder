import base64
import html
import os
import re
import sqlite3
import json
import uuid
import textwrap
from datetime import datetime
from io import BytesIO

import streamlit as st

try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    REPORTLAB_AVAILABLE = True
except Exception:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
except Exception:
    MATPLOTLIB_AVAILABLE = False

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False

st.set_page_config(page_title="GeniA Innovation Builder", page_icon="🧠", layout="wide")


# ============================================================
# PERSISTENCIA MULTIUSUARIO (SQLite)
# Cada navegador mantiene su propia sesión de Streamlit y cada
# clic en Guardar crea un registro independiente en la base.
# ============================================================
DB_PATH = os.getenv("GENIA_DB_PATH", "genia_proyectos.db")


def get_db_connection():
    conn = sqlite3.connect(DB_PATH, timeout=30, check_same_thread=False)
    conn.execute("PRAGMA journal_mode=WAL;")
    conn.execute("PRAGMA busy_timeout=30000;")
    return conn


def init_database():
    with get_db_connection() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS proyectos_mvp (
                id TEXT PRIMARY KEY,
                nombre_usuario TEXT NOT NULL,
                nombre_proyecto TEXT,
                datos_json TEXT NOT NULL,
                fecha_guardado TEXT NOT NULL
            )
            """
        )
        conn.commit()


def guardar_proyecto(nombre_usuario, datos):
    registro_id = str(uuid.uuid4())
    fecha = datetime.now().isoformat(timespec="seconds")
    with get_db_connection() as conn:
        conn.execute(
            """
            INSERT INTO proyectos_mvp
                (id, nombre_usuario, nombre_proyecto, datos_json, fecha_guardado)
            VALUES (?, ?, ?, ?, ?)
            """,
            (
                registro_id,
                nombre_usuario.strip(),
                str(datos.get("nombre", "")).strip(),
                json.dumps(datos, ensure_ascii=False),
                fecha,
            ),
        )
        conn.commit()
    return registro_id, fecha


init_database()


def image_to_base64(path: str) -> str:
    with open(path, "rb") as file:
        return base64.b64encode(file.read()).decode()


def add_background(path: str) -> None:
    encoded = image_to_base64(path)
    st.markdown(
        f"""
        <style>
        .stApp {{
            background-image: linear-gradient(rgba(255,255,255,.08), rgba(255,255,255,.16)),
                              url("data:image/png;base64,{encoded}");
            background-size: cover; background-position: center; background-attachment: fixed;
        }}
        .block-container {{
            background: linear-gradient(135deg, rgba(6,42,72,.62), rgba(14,98,130,.50), rgba(240,90,40,.22));
            border-radius: 28px; padding: 2rem 2.2rem 3rem; box-shadow: 0 12px 36px rgba(0,0,0,.25);
            border: 1px solid rgba(255,255,255,.25);
        }}
        [data-testid="stSidebar"] {{ background: rgba(255,255,255,.96); }}
        .genia-header {{ display:flex; align-items:center; gap:24px; padding:20px 26px; border-radius:26px;
            background:linear-gradient(135deg, rgba(255,255,255,.60), rgba(232,246,250,.44));
            border:1px solid rgba(11,46,74,.18); box-shadow:0 8px 24px rgba(0,0,0,.12); margin-bottom:18px; }}
        .genia-logo-img {{ width:210px; max-width:28vw; height:auto; object-fit:contain; }}
        .genia-title {{ font-size:40px; line-height:1.05; font-weight:850; color:#0B2E4A; }}
        .genia-subtitle {{ font-size:20px; color:#F05A28; font-weight:750; margin-top:6px; }}
        .genia-caption {{ font-size:14px; color:#5B667A; margin-top:4px; }}
        .block-container h1,.block-container h2,.block-container h3,.block-container p,
        .block-container label,.block-container span {{ color:#fff !important; }}
        input,textarea {{ background-color:rgba(255,255,255,.72)!important; color:#0B2E4A!important; }}
        div[data-baseweb="select"]>div {{ background-color:rgba(255,255,255,.72)!important; color:#0B2E4A!important; }}
        [data-testid="stExpander"] {{ background:rgba(255,255,255,.58); border-radius:14px; }}
        [data-testid="stExpander"] * {{ color:#0B2E4A!important; }}
        .stTabs [data-baseweb="tab"] {{ color:#fff!important; font-weight:650; }}
        .stTabs [aria-selected="true"] {{ color:#F05A28!important; background:rgba(255,255,255,.62); border-radius:10px 10px 0 0; }}
        .result-box,.maturity-card {{
            background:rgba(255,255,255,.86)!important;
            border-radius:22px;
            padding:22px;
            margin-top:12px;
            border:1px solid rgba(255,255,255,.28);
            box-shadow:0 10px 26px rgba(0,0,0,.16);
            backdrop-filter: blur(6px);
        }}
        .block-container .result-box,
        .block-container .result-box h1,
        .block-container .result-box h2,
        .block-container .result-box h3,
        .block-container .result-box h4,
        .block-container .result-box p,
        .block-container .result-box span,
        .block-container .result-box strong,
        .block-container .result-box b,
        .block-container .result-box li,
        .block-container .result-box div,
        .block-container .maturity-card,
        .block-container .maturity-card h1,
        .block-container .maturity-card h2,
        .block-container .maturity-card h3,
        .block-container .maturity-card p,
        .block-container .maturity-card span,
        .block-container .maturity-card strong,
        .block-container .maturity-card b,
        .block-container .maturity-card div {{
            color:#0B2E4A!important;
            -webkit-text-fill-color:#0B2E4A!important;
        }}
        .result-box {{
            white-space:pre-wrap;
            font-size:16px;
            line-height:1.55;
            overflow-wrap:anywhere;
        }}
        textarea {{
            background-color:rgba(255,255,255,.98)!important;
            color:#0B2E4A!important;
            -webkit-text-fill-color:#0B2E4A!important;
        }}
        .stDownloadButton button,.stButton button {{
            background:#fff!important;
            color:#0B2E4A!important;
            font-weight:700!important;
            border-radius:12px!important;
            border:1px solid rgba(11,46,74,.25)!important;
        }}
        .stDownloadButton button p,
        .stDownloadButton button span,
        .stButton button p,
        .stButton button span {{
            color:#0B2E4A!important;
            -webkit-text-fill-color:#0B2E4A!important;
        }}
        .stDownloadButton button:hover,.stButton button:hover {{
            background:#F05A28!important;
            color:#fff!important;
        }}
        .stDownloadButton button:hover p,
        .stDownloadButton button:hover span,
        .stButton button:hover p,
        .stButton button:hover span {{
            color:#fff!important;
            -webkit-text-fill-color:#fff!important;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )


if os.path.exists("fondo_bootcamp.png"):
    add_background("fondo_bootcamp.png")

if os.path.exists("logo_genia.png"):
    logo = f'<img src="data:image/png;base64,{image_to_base64("logo_genia.png")}" class="genia-logo-img">'
else:
    logo = '<div style="font-size:48px;">🧠</div>'

st.markdown(
    f"""
    <div class="genia-header"><div>{logo}</div><div>
    <div class="genia-title">Innovation Builder</div>
    <div class="genia-subtitle">Del reto clínico al MVP de IA en salud</div>
    <div class="genia-caption">Programa de Inteligencia Artificial — Los Cobos Medical Center</div>
    </div></div>
    """,
    unsafe_allow_html=True,
)

with st.expander("¿Qué es un MVP de IA en salud?", expanded=True):
    st.markdown("""
    Un **MVP de IA en salud** es la versión más simple, segura y medible de una solución que permite comprobar
    si genera valor clínico u operativo antes de invertir en un desarrollo completo.

    **Debe ser:** específico, verificable, medible, supervisado y limitado a un caso de uso.  
    **No debe ser:** una plataforma completa, una IA autónoma ni un reemplazo del profesional.
    """)

if "data" not in st.session_state:
    st.session_state.data = {}

if "usuario_registro" not in st.session_state:
    st.session_state.usuario_registro = ""

if "ultimo_guardado" not in st.session_state:
    st.session_state.ultimo_guardado = None

def save(key, value):
    st.session_state.data[key] = value

def text_ok(d, key, minimum=8):
    return len(str(d.get(key, "")).strip()) >= minimum

def safe_filename(text):
    return re.sub(r"[^A-Za-z0-9_-]+", "_", (text or "proyecto")).strip("_") or "proyecto"


tabs = st.tabs([
    "1. Problema", "2. Usuario y flujo", "3. Datos e IA", "4. Variable objetivo",
    "5. MVP", "6. Riesgos", "7. Validación y KPIs", "8. Piloto", "9. Escalamiento", "10. Resultado"
])

with tabs[0]:
    st.header("1. Definir el problema")
    c1, c2 = st.columns(2)
    with c1:
        nombre = st.text_input("Nombre del proyecto / equipo", value=st.session_state.data.get("nombre", ""))
        area = st.selectbox("Área principal", ["Urgencias","Radiología","Patología","Oncología","Consulta externa","Hospitalización","Laboratorio clínico","Farmacia","Gestión administrativa","Salud pública","Otra"])
        tipo_problema = st.selectbox("Tipo de problema", ["Diagnóstico","Pronóstico","Triage / priorización","Seguimiento","Gestión del flujo","Extracción de información","Resumen de historia clínica","Planeación de demanda","Educación","Otro"])
    with c2:
        problema = st.text_area("Problema específico", placeholder="Ejemplo: identificación tardía de pacientes con riesgo de deterioro.", height=110)
        evidencia_problema = st.text_area("¿Qué evidencia demuestra que existe?", placeholder="Indicadores, auditorías, eventos, tiempos o entrevistas.", height=90)
        consecuencia = st.text_area("¿Qué ocurre si no se resuelve?", height=80)
    for k, v in {"nombre":nombre,"area":area,"tipo_problema":tipo_problema,"problema":problema,"evidencia_problema":evidencia_problema,"consecuencia":consecuencia}.items(): save(k,v)

with tabs[1]:
    st.header("2. Usuario, decisión y flujo")
    c1, c2 = st.columns(2)
    with c1:
        usuario = st.selectbox("Usuario principal", ["Médico general","Especialista","Enfermería","Paciente","Administrativo","Radiólogo","Patólogo","Farmacéutico","Gestor clínico","Equipo de calidad","Otro"])
        momento = st.selectbox("Momento del flujo", ["Antes de la consulta","Durante la consulta","Después de la consulta","Ingreso a urgencias","Durante hospitalización","Antes del egreso","Lectura de imágenes","Junta médica","Seguimiento ambulatorio","Gestión operativa","Otro"])
        decision = st.text_area("¿Qué decisión o acción apoyará?", height=90)
    with c2:
        flujo_actual = st.text_area("Flujo actual sin IA", placeholder="Datos → revisión manual → decisión", height=85)
        flujo_mvp = st.text_area("Flujo propuesto con MVP", placeholder="Datos → MVP → resultado verificable → profesional decide", height=85)
        responsable_final = st.text_input("Responsable de la decisión final")
    for k, v in {"usuario":usuario,"momento":momento,"decision":decision,"flujo_actual":flujo_actual,"flujo_mvp":flujo_mvp,"responsable_final":responsable_final}.items(): save(k,v)

with tabs[3]:
    st.header("4. Variable objetivo y alcance")
    tarea_ia = st.selectbox("Tarea principal", ["Predecir un evento","Clasificar","Priorizar","Extraer información","PLN","Pronosticar demanda","Detectar una anomalía","Recomendar bajo supervisión"])
    c1, c2 = st.columns(2)
    with c1:
        variable_objetivo = st.text_area("Variable objetivo o resultado principal", placeholder="Ejemplo: reingreso no programado dentro de 30 días.", height=90)
        definicion_positiva = st.text_area("Definición exacta del resultado correcto", height=90)
        unidad_analisis = st.selectbox("Unidad de análisis", ["Paciente","Consulta","Hospitalización","Imagen","Documento","Medicamento","Día","Servicio","Otra"])
    with c2:
        horizonte = st.text_input("Horizonte temporal", placeholder="24 horas, 7 días, 30 días")
        poblacion = st.text_area("Población objetivo", height=85)
        exclusion = st.text_area("Casos fuera del alcance", height=85)
    for k, v in {"tarea_ia":tarea_ia,"variable_objetivo":variable_objetivo,"definicion_positiva":definicion_positiva,"unidad_analisis":unidad_analisis,"horizonte":horizonte,"poblacion":poblacion,"exclusion":exclusion}.items(): save(k,v)

with tabs[2]:
    st.header("3. Datos, línea base y técnica de IA")
    datos = st.multiselect("Datos necesarios", ["Historia clínica","Notas médicas","Signos vitales","Laboratorio","Imágenes médicas","Citología / patología digital","Medicamentos","Encuestas","Datos administrativos","Series temporales","Datos externos","Otro"])
    c1, c2 = st.columns(2)
    with c1:
        fuente = st.text_area("Fuente de los datos", height=80)
        disponibilidad = st.selectbox("Disponibilidad real", ["No confirmada","Parcial","Disponible sin depurar","Disponible y depurada","Disponible y lista para análisis"])
        volumen = st.number_input("Número aproximado de registros", min_value=0, max_value=10000000, value=0, step=10)
        calidad = st.select_slider("Calidad percibida", options=[0,1,2,3,4,5], value=0, format_func=lambda x: "No evaluada" if x == 0 else str(x))
    with c2:
        etiqueta = st.text_area("¿Cómo se obtiene la verdad de referencia?", height=80)
        linea_base = st.text_area("Línea base de comparación", placeholder="Revisión manual, regla clínica, promedio histórico o modelo simple.", height=80)
        tecnica = st.multiselect("Técnica probable", ["Reglas clínicas","Regresión logística","Machine Learning clásico","Gradient Boosting","Deep Learning","Computer Vision","NLP","IA generativa","Series temporales","Reglas + IA","No definido aún"])
    for k, v in {"datos":datos,"fuente":fuente,"disponibilidad":disponibilidad,"volumen":volumen,"calidad":calidad,"etiqueta":etiqueta,"linea_base":linea_base,"tecnica":tecnica}.items(): save(k,v)

with tabs[4]:
    st.header("5. Diseñar el MVP")
    st.info("El MVP debe hacer una sola cosa útil, verificable y segura.")
    funcion = st.text_area("Función mínima", placeholder="Mi MVP hará solamente...", height=85)
    no_hara = st.text_area("¿Qué NO hará?", placeholder="No diagnosticará, no formulará tratamiento, no reemplazará criterio profesional.", height=85)
    c1, c2 = st.columns(2)
    with c1:
        entrada_mvp = st.text_area("Entrada", height=75)
        procesamiento_mvp = st.text_area("Procesamiento mínimo", height=75)
        salida = st.selectbox("Salida principal", ["Alerta","Resumen","Priorización","Clasificación de riesgo","Pronóstico","Extracción estructurada","Recomendación supervisada","Reporte","Tablero","Otro"])
    with c2:
        salida_visible = st.text_area("¿Qué verá exactamente el usuario?", height=75)
        accion_usuario = st.text_area("¿Qué hará el usuario después?", height=75)
        valor = st.text_area("Valor clínico u operativo esperado", height=75)
    for k, v in {"funcion":funcion,"no_hara":no_hara,"entrada_mvp":entrada_mvp,"procesamiento_mvp":procesamiento_mvp,"salida":salida,"salida_visible":salida_visible,"accion_usuario":accion_usuario,"valor":valor}.items(): save(k,v)

with tabs[5]:
    st.header("6. Riesgos, supervisión y gobernanza")
    riesgo = st.selectbox("Riesgo principal", ["Falso negativo","Falso positivo","Sesgo poblacional","Sobreconfianza","Datos incompletos","Privacidad","Falla de integración","Alucinación","Uso fuera de alcance","Otro"])
    c1, c2 = st.columns(2)
    with c1:
        peor = st.text_area("Peor escenario", height=80)
        mitigacion = st.text_area("Mitigación", height=80)
        seguridad = st.slider("Nivel estimado de control", 1, 5, 3)
    with c2:
        supervisor = st.text_input("¿Quién revisará la salida?")
        baja_confianza = st.text_area("¿Qué hará el sistema con baja confianza o datos insuficientes?", height=80)
        correccion_usuario = st.selectbox("¿El usuario puede corregir o rechazar?", ["No definido","Sí","No"])
    privacidad = st.multiselect("Controles previstos", ["Anonimización / seudonimización","Control de acceso","Registro de auditoría","Consentimiento o base legal","Aprobación institucional","Comité de ética","Evaluación de seguridad","Gestión de incidentes"])
    for k, v in {"riesgo":riesgo,"peor":peor,"mitigacion":mitigacion,"seguridad":seguridad,"supervisor":supervisor,"baja_confianza":baja_confianza,"correccion_usuario":correccion_usuario,"privacidad":privacidad}.items(): save(k,v)

with tabs[6]:
    st.header("7. Validación y KPIs")
    validaciones = st.multiselect("Tipos de validación", ["Validación técnica interna","Validación con expertos","Prueba de usabilidad","Validación clínica retrospectiva","Validación prospectiva","Evaluación operativa","Evaluación de sesgo","Evaluación de calibración"])
    c1, c2 = st.columns(2)
    with c1:
        kpis = st.multiselect("Máximo 4 KPIs", ["Tiempo ahorrado","Concordancia con experto","Sensibilidad","Especificidad","Precisión","AUPRC","AUROC","F1","Calibración","Tasa de falsos negativos","Satisfacción del usuario","Casos priorizados correctamente","Reducción de errores","Costo por caso","Adherencia al uso","MAE / RMSE","WAPE / MAPE"], max_selections=4)
        meta = st.text_area("Meta cuantificable", height=80)
    with c2:
        medicion = st.text_area("Cómo se medirá", height=80)
        muestra_validacion = st.number_input("Casos para validación inicial", min_value=0, max_value=1000000, value=0, step=10)
        criterio_avance = st.text_area("Criterio para avanzar al piloto", height=80)
    for k, v in {"validaciones":validaciones,"kpis":kpis,"meta":meta,"medicion":medicion,"muestra_validacion":muestra_validacion,"criterio_avance":criterio_avance}.items(): save(k,v)

with tabs[7]:
    st.header("8. Plan piloto")
    c1, c2 = st.columns(2)
    with c1:
        lugar = st.selectbox("Lugar del piloto", ["Urgencias","Radiología","Patología","Consulta externa","Hospitalización","Farmacia","Comité / junta médica","Gestión administrativa","Otro"])
        duracion = st.selectbox("Duración", ["1 semana","2 semanas","4 semanas","8 semanas","12 semanas"])
        usuarios = st.number_input("Usuarios participantes", min_value=1, max_value=1000, value=5)
    with c2:
        casos = st.number_input("Casos a evaluar", min_value=1, max_value=100000, value=50)
        criterio_stop = st.text_area("Criterio de detención", height=80)
        responsable_piloto = st.text_input("Responsable del piloto")
    for k, v in {"lugar":lugar,"duracion":duracion,"usuarios":usuarios,"casos":casos,"criterio_stop":criterio_stop,"responsable_piloto":responsable_piloto}.items(): save(k,v)

with tabs[8]:
    st.header("9. Escalamiento y propiedad intelectual")
    escala = st.text_area("Ruta de escalamiento", height=90)
    pi = st.multiselect("Propiedad intelectual potencial", ["Software","Algoritmo / modelo","Dataset curado","Flujo clínico","Interfaz","Método de validación","Marca","Secreto industrial"])
    aliados = st.text_area("Aliados necesarios", height=75)
    integracion = st.multiselect("Integraciones futuras", ["HIS","PACS","LIS","ERP","API","Repositorio documental","Otra"])
    for k, v in {"escala":escala,"pi":pi,"aliados":aliados,"integracion":integracion}.items(): save(k,v)


def maturity_score(d):
    score = 0
    weighted = {
        "problema":4,"evidencia_problema":3,"decision":3,"flujo_mvp":3,"responsable_final":3,
        "variable_objetivo":5,"definicion_positiva":4,"horizonte":3,"poblacion":2,"exclusion":2,
        "fuente":3,"etiqueta":3,"linea_base":4,"funcion":4,"no_hara":3,"entrada_mvp":2,
        "salida_visible":2,"accion_usuario":3,"valor":2,"peor":3,"mitigacion":4,"supervisor":2,
        "baja_confianza":3,"meta":3,"criterio_avance":3,"criterio_stop":2,"escala":2
    }
    for key, points in weighted.items():
        if text_ok(d, key, 6 if points <= 2 else 8): score += points
    score += 3 if d.get("datos") else 0
    score += 4 if d.get("disponibilidad") in ["Disponible sin depurar","Disponible y depurada","Disponible y lista para análisis"] else 0
    score += 2 if int(d.get("volumen",0)) > 0 else 0
    score += 2 if d.get("correccion_usuario") == "Sí" else 0
    score += 2 if len(d.get("privacidad",[])) >= 2 else 0
    score += 3 if len(d.get("validaciones",[])) >= 2 else 0
    score += 3 if len(d.get("kpis",[])) >= 2 else 0
    score += 2 if int(d.get("muestra_validacion",0)) > 0 else 0
    return min(score,100)


def maturity_level(score, d):
    blockers = []
    for key, label in [("variable_objetivo","variable objetivo"),("linea_base","línea base"),("mitigacion","mitigación"),("criterio_avance","criterio de avance")]:
        if not text_ok(d,key,8): blockers.append(label)
    if len(d.get("validaciones",[])) < 2: blockers.append("plan de validación")
    if score < 40: return "🔴 NIVEL 1: IDEA", "El problema está identificado, pero el MVP aún no está definido.", "Definir objetivo, datos, función mínima y riesgos.", blockers
    if score < 65: return "🟡 NIVEL 2: PROTOTIPO CONCEPTUAL", "Existe una estructura inicial, pero falta viabilidad.", "Fortalecer datos, línea base, validación y supervisión.", blockers
    if score < 85 or blockers: return "🟢 NIVEL 3: MVP DEFINIDO", "El MVP puede avanzar a validación inicial.", "Resolver elementos bloqueantes antes del piloto.", blockers
    return "🔵 NIVEL 4: CANDIDATO A PILOTO", "Cuenta con objetivo, datos, seguridad y validación.", "Realizar revisión institucional antes de ejecutarlo.", blockers


def innovation_score(d):
    return {
        "Impacto": min(10, (4 if text_ok(d,"problema",20) else 0)+(3 if text_ok(d,"valor") else 0)+(3 if text_ok(d,"meta") else 0)),
        "Factibilidad": min(10, (3 if text_ok(d,"funcion") else 0)+(3 if text_ok(d,"linea_base") else 0)+(4 if d.get("disponibilidad") in ["Disponible sin depurar","Disponible y depurada","Disponible y lista para análisis"] else 0)),
        "Datos": min(10, int(d.get("calidad",0))*2 + (2 if text_ok(d,"etiqueta") else 0)),
        "Escalabilidad": min(10, (4 if text_ok(d,"escala") else 0)+(3 if text_ok(d,"aliados") else 0)+(2 if d.get("pi") else 0)+(1 if d.get("integracion") else 0)),
        "Control del riesgo": min(10, (3 if text_ok(d,"mitigacion") else 0)+(2 if text_ok(d,"baja_confianza") else 0)+(2 if text_ok(d,"supervisor") else 0)+(2 if len(d.get("privacidad",[]))>=2 else 0)+(1 if d.get("correccion_usuario")=="Sí" else 0))
    }


def recommendations(d):
    recs=[]
    checks=[("variable_objetivo","Definir una variable objetivo concreta y medible."),("definicion_positiva","Precisar cómo se determina el resultado correcto."),("linea_base","Definir una línea base antes de comparar modelos avanzados."),("baja_confianza","Diseñar el comportamiento ante baja confianza o datos incompletos."),("criterio_avance","Establecer un criterio cuantificable para avanzar al piloto."),("criterio_stop","Definir cuándo detener el piloto por seguridad.")]
    for key,msg in checks:
        if not text_ok(d,key,8): recs.append(msg)
    if d.get("disponibilidad") in ["No confirmada","Parcial"]: recs.append("Confirmar acceso, permisos, calidad y volumen real de datos.")
    if len(d.get("validaciones",[])) < 2: recs.append("Incluir validación técnica y validación con expertos o usuarios.")
    if d.get("correccion_usuario") != "Sí": recs.append("Permitir que el profesional corrija o rechace la salida.")
    return recs[:7] or ["El proyecto está bien estructurado; sigue la revisión institucional."]


def build_summary(d, score, level, innovation, blockers):
    return f"""# {d.get('nombre','MVP sin nombre')}

## Problema
**Área:** {d.get('area','')}  
**Problema:** {d.get('problema','')}  
**Evidencia:** {d.get('evidencia_problema','')}  
**Usuario:** {d.get('usuario','')}  
**Decisión:** {d.get('decision','')}

## Variable objetivo
**Tarea:** {d.get('tarea_ia','')}  
**Resultado:** {d.get('variable_objetivo','')}  
**Definición:** {d.get('definicion_positiva','')}  
**Horizonte:** {d.get('horizonte','')}  
**Población:** {d.get('poblacion','')}

## Datos y línea base
**Datos:** {', '.join(d.get('datos',[]))}  
**Fuente:** {d.get('fuente','')}  
**Disponibilidad:** {d.get('disponibilidad','')}  
**Verdad de referencia:** {d.get('etiqueta','')}  
**Línea base:** {d.get('linea_base','')}

## MVP
**Función mínima:** {d.get('funcion','')}  
**No hará:** {d.get('no_hara','')}  
**Entrada:** {d.get('entrada_mvp','')}  
**Salida:** {d.get('salida_visible','')}  
**Acción del usuario:** {d.get('accion_usuario','')}

## Seguridad
**Riesgo:** {d.get('riesgo','')}  
**Mitigación:** {d.get('mitigacion','')}  
**Supervisor:** {d.get('supervisor','')}  
**Baja confianza:** {d.get('baja_confianza','')}

## Validación y piloto
**Validaciones:** {', '.join(d.get('validaciones',[]))}  
**KPIs:** {', '.join(d.get('kpis',[]))}  
**Meta:** {d.get('meta','')}  
**Piloto:** {d.get('lugar','')} durante {d.get('duracion','')} con {d.get('casos','')} casos

## Resultado automático
**Madurez:** {level} ({score}/100)  
**Innovación:** {innovation}/50  
**Bloqueantes:** {', '.join(blockers) if blockers else 'Ninguno'}

## Recomendaciones
""" + "\n".join(f"- {r}" for r in recommendations(d))


def make_pdf(d, summary):
    if not REPORTLAB_AVAILABLE: return None
    buffer=BytesIO(); doc=SimpleDocTemplate(buffer,pagesize=letter,rightMargin=36,leftMargin=36,topMargin=36,bottomMargin=36)
    styles=getSampleStyleSheet(); story=[Paragraph("GeniA Innovation Builder",styles["Title"]),Spacer(1,12)]
    rows=[["Campo","Respuesta"],["Proyecto",d.get("nombre","")],["Problema",d.get("problema","")],["Usuario",d.get("usuario","")],["Variable objetivo",d.get("variable_objetivo","")],["Función mínima",d.get("funcion","")],["Línea base",d.get("linea_base","")],["Riesgo",d.get("riesgo","")],["Mitigación",d.get("mitigacion","")],["KPIs",", ".join(d.get("kpis",[]))],["Piloto",f"{d.get('lugar','')} | {d.get('duracion','')} | {d.get('casos','')} casos"]]
    table=Table(rows,colWidths=[125,365],repeatRows=1)
    table.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.HexColor("#0B2E4A")),("TEXTCOLOR",(0,0),(-1,0),colors.white),("GRID",(0,0),(-1,-1),.5,colors.grey),("VALIGN",(0,0),(-1,-1),"TOP"),("BACKGROUND",(0,1),(0,-1),colors.HexColor("#EAF3F8"))]))
    story += [table,Spacer(1,12),Paragraph(f"Fecha: {datetime.now().strftime('%Y-%m-%d %H:%M')}",styles["Normal"])]
    doc.build(story); buffer.seek(0); return buffer


def build_pitch_slides(d, score, level, innovation):
    return [
        {
            "title": "1. El problema",
            "items": [
                f"Proyecto: {d.get('nombre','') or 'Por definir'}",
                f"Área: {d.get('area','') or 'Por definir'}",
                f"Problema: {d.get('problema','') or 'Por definir'}",
                f"Evidencia: {d.get('evidencia_problema','') or 'Por definir'}",
                f"Consecuencia: {d.get('consecuencia','') or 'Por definir'}",
            ],
        },
        {
            "title": "2. Usuario y oportunidad",
            "items": [
                f"Usuario principal: {d.get('usuario','') or 'Por definir'}",
                f"Momento del flujo: {d.get('momento','') or 'Por definir'}",
                f"Decisión apoyada: {d.get('decision','') or 'Por definir'}",
                f"Responsable final: {d.get('responsable_final','') or 'Por definir'}",
                f"Población objetivo: {d.get('poblacion','') or 'Por definir'}",
            ],
        },
        {
            "title": "3. La solución MVP",
            "items": [
                f"Función mínima: {d.get('funcion','') or 'Por definir'}",
                f"Variable objetivo: {d.get('variable_objetivo','') or 'Por definir'}",
                f"Entrada: {d.get('entrada_mvp','') or 'Por definir'}",
                f"Salida: {d.get('salida_visible','') or d.get('salida','') or 'Por definir'}",
                f"Acción del usuario: {d.get('accion_usuario','') or 'Por definir'}",
                f"Línea base: {d.get('linea_base','') or 'Por definir'}",
            ],
        },
        {
            "title": "4. Evidencia, seguridad y métricas",
            "items": [
                f"Datos: {', '.join(d.get('datos',[])) or 'Por definir'}",
                f"Técnica probable: {', '.join(d.get('tecnica',[])) or 'Por definir'}",
                f"Riesgo principal: {d.get('riesgo','') or 'Por definir'}",
                f"Mitigación: {d.get('mitigacion','') or 'Por definir'}",
                f"KPIs: {', '.join(d.get('kpis',[])) or 'Por definir'}",
                f"Meta: {d.get('meta','') or 'Por definir'}",
            ],
        },
        {
            "title": "5. Piloto y llamado a la acción",
            "items": [
                f"Piloto: {d.get('lugar','') or 'Por definir'} | {d.get('duracion','') or 'Por definir'} | {d.get('casos','')} casos",
                f"Responsable: {d.get('responsable_piloto','') or 'Por definir'}",
                f"Escalamiento: {d.get('escala','') or 'Por definir'}",
                f"Madurez: {level} ({score}/100)",
                f"Innovación: {innovation}/50",
                "Siguiente paso: validar, priorizar y preparar el piloto institucional.",
            ],
        },
    ]


def make_pptx(slides):
    if not PPTX_AVAILABLE:
        return None
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 46, 74)
        title_box = slide.shapes.add_textbox(Inches(.7), Inches(.45), Inches(12), Inches(.8))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = slide_data["title"]
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(240, 90, 40)
        body_box = slide.shapes.add_textbox(Inches(.9), Inches(1.45), Inches(11.6), Inches(5.4))
        tf = body_box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(slide_data["items"]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(255,255,255)
            p.space_after = Pt(10)
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def make_scores_chart(scores):
    if not MATPLOTLIB_AVAILABLE:
        return None

    labels = list(scores.keys())
    values = list(scores.values())

    fig, ax = plt.subplots(figsize=(6.2, 3.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#F8FBFD")

    bars = ax.bar(labels, values, width=0.58, edgecolor="#0B2E4A", linewidth=0.8)

    ax.set_ylim(0, 10)
    ax.set_ylabel("Puntaje")
    ax.set_title("Potencial de innovación", fontsize=16, pad=12)
    ax.tick_params(axis="x", rotation=18, labelsize=10)
    ax.tick_params(axis="y", labelsize=10)
    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    for bar, value in zip(bars, values):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            value + 0.18,
            str(value),
            ha="center",
            va="bottom",
            fontsize=9
        )

    fig.tight_layout()
    output = BytesIO()
    fig.savefig(output, format="png", dpi=220, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    output.seek(0)
    return output



def get_font(size=24, bold=False):
    possible_fonts = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/dejavu/DejaVuSans.ttf",
        "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "arialbd.ttf" if bold else "arial.ttf",
    ]
    for font_path in possible_fonts:
        try:
            return ImageFont.truetype(font_path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def make_canvas_png(d, score, level, innovation, blockers):
    if not PIL_AVAILABLE:
        return None

    width, height = 1800, 1250
    image = Image.new("RGB", (width, height), (11, 46, 74))
    draw = ImageDraw.Draw(image)

    title_font = get_font(44, bold=True)
    subtitle_font = get_font(28, bold=True)
    section_font = get_font(26, bold=True)
    text_font = get_font(24, bold=False)
    footer_font = get_font(22, bold=True)

    white = (255, 255, 255)
    orange = (240, 90, 40)
    dark_text = (8, 30, 48)
    light_box = (250, 250, 250)

    draw.text((55, 28), "GeniA Innovation Builder - MVP Canvas", fill=white, font=title_font)
    draw.text(
        (55, 82),
        f"Proyecto: {d.get('nombre', '') or 'Sin nombre'}",
        fill=orange,
        font=subtitle_font
    )

    def clean_value(value, fallback="Por definir"):
        if value is None:
            return fallback
        text = str(value).strip()
        return text if text else fallback

    def wrap_text(value, width_chars=48):
        lines = []
        for paragraph in str(value).split("\n"):
            wrapped = textwrap.wrap(paragraph, width=width_chars) if paragraph.strip() else [""]
            lines.extend(wrapped)
        return lines

    boxes = [
        ("Problema", clean_value(d.get("problema"))),
        (
            "Usuario y decisión",
            f"{clean_value(d.get('usuario'))}\n{clean_value(d.get('decision'))}"
        ),
        ("Variable objetivo", clean_value(d.get("variable_objetivo"))),
        (
            "Datos y línea base",
            f"{clean_value(', '.join(d.get('datos', [])), 'Sin datos definidos')}\n"
            f"Línea base: {clean_value(d.get('linea_base'))}"
        ),
        (
            "MVP",
            f"{clean_value(d.get('funcion'))}\n"
            f"Salida: {clean_value(d.get('salida_visible'))}"
        ),
        (
            "Riesgo y mitigación",
            f"{clean_value(d.get('riesgo'))}\n{clean_value(d.get('mitigacion'))}"
        ),
        (
            "KPIs y meta",
            f"{clean_value(', '.join(d.get('kpis', [])), 'Sin KPIs definidos')}\n"
            f"Meta: {clean_value(d.get('meta'))}"
        ),
        (
            "Piloto y escalamiento",
            f"{clean_value(d.get('lugar'))} | {clean_value(d.get('duracion'))} | "
            f"{clean_value(d.get('casos'))} casos\n{clean_value(d.get('escala'))}"
        ),
    ]

    coords = [
        (55, 140, 860, 355),
        (935, 140, 1740, 355),
        (55, 390, 860, 605),
        (935, 390, 1740, 605),
        (55, 640, 860, 855),
        (935, 640, 1740, 855),
        (55, 890, 860, 1105),
        (935, 890, 1740, 1105),
    ]

    for (title, content), (x1, y1, x2, y2) in zip(boxes, coords):
        draw.rounded_rectangle(
            (x1, y1, x2, y2),
            radius=24,
            fill=light_box,
            outline=orange,
            width=4
        )
        draw.text((x1 + 22, y1 + 18), title, fill=dark_text, font=section_font)

        y_text = y1 + 68
        for line in wrap_text(content, width_chars=48):
            if y_text > y2 - 34:
                draw.text((x1 + 22, y_text), "…", fill=dark_text, font=text_font)
                break
            draw.text((x1 + 22, y_text), line, fill=dark_text, font=text_font)
            y_text += 32

    footer_y1, footer_y2 = 1135, 1210
    draw.rounded_rectangle((55, footer_y1, 1740, footer_y2), radius=20, fill=orange)

    blockers_text = ", ".join(blockers) if blockers else "Ninguno"
    footer_text = (
        f"Madurez: {level} ({score}/100) | Innovación: {innovation}/50 | "
        f"Bloqueantes: {blockers_text}"
    )

    footer_lines = textwrap.wrap(footer_text, width=115)
    footer_text_y = footer_y1 + 14
    for line in footer_lines[:2]:
        draw.text((75, footer_text_y), line, fill=white, font=footer_font)
        footer_text_y += 28

    buffer = BytesIO()
    image.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


with tabs[9]:
    st.header("10. Resultado del MVP")
    d=st.session_state.data; score=maturity_score(d); level, interpretation, next_step, blockers=maturity_level(score,d)
    scores=innovation_score(d); innovation=sum(scores.values())
    c1,c2,c3,c4=st.columns(4); c1.metric("Madurez",f"{score}/100"); c2.metric("Innovación",f"{innovation}/50"); c3.metric("KPIs",len(d.get("kpis",[]))); c4.metric("Validaciones",len(d.get("validaciones",[])))
    st.markdown(f'<div class="maturity-card"><h2>{level}</h2><p><b>Interpretación:</b> {interpretation}</p><p><b>Próximo paso:</b> {next_step}</p></div>',unsafe_allow_html=True)
    if blockers: st.warning("Elementos bloqueantes: " + ", ".join(blockers))
    st.subheader("Potencial de innovación")
    for item,value in scores.items(): st.progress(value/10,text=f"{item}: {value}/10")
    st.subheader("Recomendaciones")
    for rec in recommendations(d): st.write("• "+rec)
    summary=build_summary(d,score,level,innovation,blockers)
    st.subheader("Resumen del MVP")
    summary_html = html.escape(summary)
    st.markdown(
        f'<div class="result-box">{summary_html}</div>',
        unsafe_allow_html=True
    )

    filename=safe_filename(d.get("nombre"))
    st.download_button(
        "📥 Descargar resumen en Markdown",
        summary.encode("utf-8"),
        file_name=f"genia_mvp_{filename}.md",
        mime="text/markdown"
    )

    pitch_slides = build_pitch_slides(d, score, level, innovation)
    st.subheader("Elevator pitch en 5 diapositivas")
    st.caption("Una idea central por diapositiva. Presentación sugerida: 3 a 5 minutos. Las tarjetas tienen un fondo translúcido para facilitar la lectura.")
    for slide_data in pitch_slides:
        title_html = html.escape(slide_data["title"])
        content_html = "<br>".join("• " + html.escape(item) for item in slide_data["items"])
        st.markdown(
            f"<div class='result-box'><h3>{title_html}</h3><div>{content_html}</div></div>",
            unsafe_allow_html=True
        )

    pitch_md = "# Elevator pitch en 5 diapositivas\n\n" + "\n\n---\n\n".join(
        f"## {slide_data['title']}\n\n" + "\n".join(f"- {item}" for item in slide_data["items"])
        for slide_data in pitch_slides
    )
    st.download_button(
        "🎤 Descargar elevator pitch en Markdown",
        pitch_md.encode("utf-8"),
        file_name=f"elevator_pitch_5_slides_{filename}.md",
        mime="text/markdown"
    )

    pptx_file = make_pptx(pitch_slides)
    if pptx_file:
        st.download_button(
            "📊 Descargar elevator pitch en PowerPoint (.pptx)",
            pptx_file,
            file_name=f"pitch_5_slides_{filename}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    else:
        st.caption("Para exportar PowerPoint instala python-pptx.")

    st.subheader("Gráfica de indicadores")
    chart_file = make_scores_chart(scores)
    if chart_file:
        chart_bytes = chart_file.getvalue()
        c_left, c_mid, c_right = st.columns([0.16, 0.68, 0.16])
        with c_mid:
            st.image(
                chart_bytes,
                caption="Vista previa de los indicadores del proyecto",
                use_container_width=True
            )
        st.download_button(
            "📈 Descargar gráfica de indicadores (.png)",
            chart_bytes,
            file_name=f"indicadores_{filename}.png",
            mime="image/png"
        )
    else:
        st.caption("Para mostrar y exportar la gráfica instala matplotlib.")

    st.subheader("Canvas visual del MVP")
    canvas_file = make_canvas_png(d, score, level, innovation, blockers)
    if canvas_file:
        canvas_bytes = canvas_file.getvalue()
        c_left, c_mid, c_right = st.columns([0.08, 0.84, 0.08])
        with c_mid:
            st.image(
                canvas_bytes,
                caption="Vista previa del Canvas del MVP",
                use_container_width=True
            )
        st.download_button(
            "🖼️ Descargar Canvas visual (.png)",
            canvas_bytes,
            file_name=f"canvas_mvp_{filename}.png",
            mime="image/png"
        )
    else:
        st.caption("Para mostrar y exportar el Canvas visual instala pillow.")

    pdf=make_pdf(d,summary)
    if pdf: st.download_button("📄 Descargar Canvas en PDF",pdf,file_name=f"genia_mvp_canvas_{filename}.pdf",mime="application/pdf")
    else: st.caption("Para exportar PDF instala reportlab: pip install reportlab")

st.sidebar.title("Guía docente")
st.sidebar.markdown("""
1. Cada grupo completa las secciones 1–9.  
2. Revisa madurez y bloqueantes.  
3. Ajusta su propuesta.  
4. Descarga el Canvas.  
5. Presenta un elevator pitch de 5 diapositivas.

**Principio:** el algoritmo no es el MVP. El MVP integra problema, usuario, datos, salida, acción, validación y seguridad.
""")
st.sidebar.markdown("**Dependencias:** streamlit, reportlab, python-pptx, pillow y matplotlib")

st.sidebar.divider()
st.sidebar.subheader("Guardar proyecto")
st.sidebar.caption("Cada usuario y dispositivo guarda un registro independiente.")
nombre_usuario_guardado = st.sidebar.text_input(
    "Nombre del participante / equipo",
    value=st.session_state.usuario_registro,
)
st.session_state.usuario_registro = nombre_usuario_guardado

if st.sidebar.button("💾 Guardar proyecto", use_container_width=True, type="primary"):
    nombre_usuario = st.session_state.usuario_registro.strip()
    if not nombre_usuario:
        st.sidebar.error("Escribe el nombre del participante o equipo antes de guardar.")
    else:
        try:
            registro_id, fecha = guardar_proyecto(nombre_usuario, st.session_state.data)
            st.session_state.ultimo_guardado = {
                "id": registro_id,
                "fecha": fecha,
                "usuario": nombre_usuario,
            }
            st.sidebar.success(f"Proyecto guardado correctamente: {fecha}")
        except sqlite3.Error as exc:
            st.sidebar.error(f"No fue posible guardar el proyecto: {exc}")

if st.session_state.ultimo_guardado:
    st.sidebar.caption(
        f"Último registro: {st.session_state.ultimo_guardado['usuario']} · "
        f"{st.session_state.ultimo_guardado['fecha']}"
    )

if st.sidebar.button("Reiniciar formulario"):
    st.session_state.data={}; st.rerun()
